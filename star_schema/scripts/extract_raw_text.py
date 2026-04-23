#!/usr/bin/env python3
"""
extract_raw_text.py — Techpack 原始資料統一提取 (STEP 1)

三條提取路徑，全部以 D-number 為 key：
  A. PDF → metadata (design_id, department, category, season, brand...)
  B. PPTX → TXT (中文翻譯做工描述)
  C. PDF → detect construction pages → render PNG (給 VLM 辨識做工)

Usage:
  python extract_raw_text.py --scan-dir ../../2026               # 跑全部
  python extract_raw_text.py --scan-dir ../../2026 --pptx-only   # 只跑 PPTX
  python extract_raw_text.py --scan-dir ../../2026 --pdf-only    # 只跑 PDF
  python extract_raw_text.py --scan-dir ../../2026 --dry-run     # 列出待處理，不實際跑
  python extract_raw_text.py --scan-dir ../../2026 --force       # 忽略已處理，全部重跑

Output (寫到 --output-dir，預設 star_schema/data/ingest/):
  metadata/designs.jsonl           每 design 一行，PDF 提取的 metadata
  pptx/facts_raw/{DID}/*.txt       每個 PPTX slide 的原始文字
  pdf/callout_images/{DID}_p{N}.png  construction page 渲染圖
  pdf/callout_manifest.jsonl       PNG 清單（append）
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
from pathlib import Path
from datetime import datetime, timezone


# ════════════════════════════════════════════════════════════
# COMMON
# ════════════════════════════════════════════════════════════

def extract_design_id_from_path(path: str) -> str | None:
    """Extract D-number from filename or path."""
    m = re.search(r"\b(D\d{3,6})\b", path)
    return m.group(1) if m else None


def extract_design_id_from_pdf_content(pdf_path: str) -> str | None:
    """Extract D-number by reading PDF first page (for ONY_ filenames)."""
    import fitz
    try:
        doc = fitz.open(pdf_path)
        if doc.page_count > 0:
            text = doc[0].get_text()[:600]
            m = re.search(r"\b(D\d{3,6})\b", text)
            doc.close()
            return m.group(1) if m else None
        doc.close()
    except Exception:
        pass
    return None


def resolve_design_id(file_path: str, is_pdf: bool = False) -> str | None:
    """Get D-number: try filename first, then PDF content if needed."""
    did = extract_design_id_from_path(file_path)
    if did:
        return did
    if is_pdf:
        return extract_design_id_from_pdf_content(file_path)
    return None


def scan_files(scan_dir: str, ext: str) -> list[str]:
    """Recursively find all files with given extension."""
    found = []
    for root, dirs, files in os.walk(scan_dir):
        for f in files:
            if f.endswith(ext) and not f.startswith("~$"):
                found.append(os.path.join(root, f))
    return sorted(found)


# ════════════════════════════════════════════════════════════
# PATH A: PDF → METADATA
# ════════════════════════════════════════════════════════════

# Centric 8 PDF first-page fields (key-value pairs in fixed layout)
CENTRIC_FIELDS = [
    "Season", "Brand/Division", "Department", "Collection",
    "Category", "Sub- Category", "Sub-Category",
    "Design Type", "Design Sub- Type", "Design Sub-Type",
    "Fit Camp", "Sleeve Length", "Rise", "Body Length", "Leg Shape",
    "Flow", "Vendor", "Status", "BOM Number", "Carry Over",
]

# Fields that signal "stop reading value lines"
STOP_FIELDS = set(CENTRIC_FIELDS) | {
    "Image", "Tech Pack", "Created", "Modified", "Modiﬁed",
    "Design BOM", "Legacy Style Numbers", "This Gap Inc.",
}


def extract_pdf_metadata(pdf_path: str) -> dict | None:
    """Extract structured metadata from Centric 8 PDF first page.

    Returns dict with: design_id, design_name, season, brand_division,
    department, collection, category, sub_category, design_type,
    fit_camp, status, bom_number, vendor, source_file, ...
    """
    import fitz

    try:
        doc = fitz.open(pdf_path)
    except Exception:
        return None

    if doc.page_count == 0:
        doc.close()
        return None

    # Read first 2 pages (some PDFs split metadata across pages)
    full_text = ""
    for i in range(min(2, doc.page_count)):
        full_text += doc[i].get_text() + "\n"
    doc.close()

    # Extract D-number
    m = re.search(r"\b(D\d{3,6})\b", full_text[:600])
    if not m:
        return None
    design_id = m.group(1)

    # Extract design name (pattern: D##### <name> ########## <status>)
    name_m = re.search(
        rf"{design_id}[- ]+(.+?)(?:\s+\d{{9}}|\s+Adopted|\s+IN WORK|\s+Concept)",
        full_text[:800],
    )
    design_name = name_m.group(1).strip() if name_m else ""

    # Parse key-value fields from Centric 8 layout
    lines = [l.strip() for l in full_text.split("\n") if l.strip()]
    fields = {}
    for i, line in enumerate(lines):
        for field in CENTRIC_FIELDS:
            if line == field:
                # Collect value from subsequent lines until next known field
                val_parts = []
                for j in range(i + 1, min(i + 5, len(lines))):
                    if lines[j] in STOP_FIELDS:
                        break
                    val_parts.append(lines[j])
                if val_parts:
                    # Normalize field name
                    key = field.replace(" ", "_").replace("-", "_").replace("/", "_").lower()
                    fields[key] = " ".join(val_parts)
                break

    # Build metadata record
    meta = {
        "design_id": design_id,
        "design_name": design_name,
        "season": fields.get("season", ""),
        "brand_division": fields.get("brand_division", ""),
        "department": fields.get("department", ""),
        "collection": fields.get("collection", ""),
        "category": fields.get("category", ""),
        "sub_category": fields.get("sub__category", "") or fields.get("sub_category", ""),
        "design_type": fields.get("design_type", ""),
        "design_sub_type": fields.get("design_sub__type", "") or fields.get("design_sub_type", ""),
        "fit_camp": fields.get("fit_camp", ""),
        "rise": fields.get("rise", ""),
        "status": fields.get("status", ""),
        "bom_number": fields.get("bom_number", ""),
        "vendor": fields.get("vendor", ""),
        "source_file": os.path.basename(pdf_path),
        "total_pages": 0,
    }

    # Re-open to get page count (cheap)
    try:
        import fitz
        doc = fitz.open(pdf_path)
        meta["total_pages"] = doc.page_count
        doc.close()
    except Exception:
        pass

    return meta


# ════════════════════════════════════════════════════════════
# PATH B: PPTX → TXT
# ════════════════════════════════════════════════════════════

def extract_pptx_text(pptx_path: str) -> dict:
    """Extract all text from PPTX. Returns {slides, total_lines, total_chars}."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides = []
    total_lines = 0
    total_chars = 0

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        slides.append({"index": i + 1, "texts": texts})
        total_lines += len(texts)
        total_chars += sum(len(t) for t in texts)

    return {"slides": slides, "total_lines": total_lines, "total_chars": total_chars}


def pptx_to_txt(pptx_path: str, out_path: str) -> dict:
    """Extract PPTX → TXT file."""
    result = extract_pptx_text(pptx_path)

    lines = [
        f"SOURCE FILE: {os.path.basename(pptx_path)}",
        f"FULL PATH: {pptx_path}",
        "",
    ]
    for slide in result["slides"]:
        for t in slide["texts"]:
            lines.append(t)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    return {
        "source": pptx_path, "output": out_path,
        "slides": len(result["slides"]),
        "lines": result["total_lines"], "chars": result["total_chars"],
    }


def make_pptx_output_name(pptx_path: str, scan_dir: str) -> str:
    """Generate output TXT filename from PPTX path."""
    rel = os.path.relpath(pptx_path, scan_dir)
    name = rel.replace(os.sep, "_").replace("/", "_")
    name = re.sub(r"[^\w\-.]", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    return name.rsplit(".pptx", 1)[0] + ".txt"


# ════════════════════════════════════════════════════════════
# PATH C: PDF → CONSTRUCTION PAGE DETECT + RENDER PNG
# Based on construction-page-rules.md
# ════════════════════════════════════════════════════════════

def detect_construction_pages(pdf_path: str) -> list[dict]:
    """Detect construction callout pages in a Centric 8 PDF.

    Scoring system from construction-page-rules.md:

    Positive:
      +3  ISO code (301/401/406/514/602/605/607)
      +3  title "CONSTRUCTION CALLOUTS" / "INTERNAL/CONSTRUCTION"
      +3  "Type" + "Construction Callout" (Centric 8 metadata page)
      +2  sewing keyword
      +2  margin spec (1/4", 3/8" etc.)
      +2  in callout section ("CALLOUT" / "BOM REVIEW")
      +1  needle count (2N, 3N5TH etc.)

    Negative (hard exclude):
      GRADE REVIEW, REF IMAGES, INSPIRATION, FIT COMMENTS, PATTERN CORRECTIONS
      POM table (POM NAME / TOL FRACTION / VENDOR ACTUAL)
      Measurement Chart as primary content (with TARGET/TOLERANCE/GRADING)

    Soft negative:
      -3  ADDITIONAL COMMENTS (factory suggestion)
      -2  BOM material table (FABRIC + TRIM + SUPPLIER ≥3)

    Pass 2: adjacent low-word-count pages near confirmed callout pages
    (Centric 8 pattern: metadata page → actual callout image page)

    Returns: [{"page": 1-indexed, "type": "image"|"text", "score": N, "word_count": N}]
    """
    import fitz

    doc = fitz.open(pdf_path)
    results = []

    ISO_RE = re.compile(r"\b(301|401|406|514|602|605|607)\b")
    MARGIN_RE = re.compile(r'\d+/\d+["\u201D]|1["\u201D]|1\s+1/\d+["\u201D]')
    NEEDLE_RE = re.compile(r"\b[23]N\b|\b[23]NDL\b|\b2N3TH\b|\b3N5TH\b", re.I)

    SEWING_KW = [
        "COVERSTITCH", "OVERLOCK", "TOPSTITCH", "FLATLOCK", "FLATSEAM",
        "BARTACK", "BAR TACK", "BLINDHEM", "BLIND HEM", "EDGESTITCH",
        "EDGE STITCH", "CHAINSTITCH", "CHAIN STITCH", "FELLED SEAM",
        "LAPPED SEAM", "SATIN STITCH", "CLEAN FINISH", "TURNBACK",
        "TURN BACK", "UNDERSTITCHED", "BINDING", "SERGE", "SERGED",
    ]

    EXCLUDE_TITLES = [
        "GRADE REVIEW", "REF IMAGES", "REFERENCE IMAGES", "INSPIRATION",
        "FIT COMMENTS", "FIT SAMPLE IMAGES", "PATTERN CORRECTIONS",
        "NEXT STEPS", "FITSAMPLEIMAGES", "FITCOMMENTS",
    ]

    POM_KW = ["POM NAME", "TOL FRACTION", "VENDOR ACTUAL", "SAMPLE EVAL",
              "QC EVALUATION"]

    # Pass 1: score each page
    for i in range(doc.page_count):
        page = doc[i]
        text = page.get_text()
        text_upper = text.upper()
        word_count = len(text.split())

        if any(kw in text_upper for kw in EXCLUDE_TITLES):
            continue
        if any(kw in text_upper for kw in POM_KW):
            continue
        if "MEASUREMENT CHART" in text_upper:
            if any(kw in text_upper for kw in ["TARGET", "TOLERANCE", "GRADING"]):
                continue

        score = 0
        if ISO_RE.search(text):                                          score += 3
        if "CONSTRUCTION CALLOUT" in text_upper or \
           "INTERNAL/CONSTRUCTION" in text_upper:                        score += 3
        if any(kw in text_upper for kw in
               ["CALLOUT", "BOM REVIEW", "DESIGN BOM"]):                score += 2
        if any(kw in text_upper for kw in SEWING_KW):                   score += 2
        if MARGIN_RE.search(text):                                       score += 2
        if NEEDLE_RE.search(text):                                       score += 1
        if "ADDITIONAL COMMENTS" in text_upper:                          score -= 3
        bom_hits = sum(1 for kw in ["FABRIC", "TRIM", "SUPPLIER", "THREAD"]
                       if kw in text_upper)
        if bom_hits >= 3:                                                score -= 2
        if "CONSTRUCTION CALLOUT" in text_upper and "TYPE" in text_upper: score += 3

        if score < 5:
            continue

        results.append({
            "page": i + 1,
            "type": "image" if word_count < 40 else "text",
            "score": score,
            "word_count": word_count,
        })

    # Pass 2: adjacent pages near confirmed callouts (Centric 8 pattern)
    confirmed = {r["page"] for r in results}
    if confirmed:
        for i in range(doc.page_count):
            pg = i + 1
            if pg in confirmed:
                continue
            if not any(abs(pg - cp) <= 2 for cp in confirmed):
                continue
            text = doc[i].get_text()
            text_upper = text.upper()
            wc = len(text.split())
            if wc < 60:
                if any(kw in text_upper for kw in EXCLUDE_TITLES + POM_KW):
                    continue
                results.append({"page": pg, "type": "image", "score": 5, "word_count": wc})

    results.sort(key=lambda x: x["page"])
    doc.close()
    return results


def render_pdf_pages(pdf_path: str, pages: list[dict],
                     output_dir: str, design_id: str) -> list[dict]:
    """Render PDF pages as 216 DPI PNG for VLM."""
    import fitz

    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    rendered = []

    for pg in pages:
        page = doc[pg["page"] - 1]
        mat = fitz.Matrix(3, 3)
        pix = page.get_pixmap(matrix=mat)
        out_name = f"{design_id}_p{pg['page']}.png"
        out_path = os.path.join(output_dir, out_name)
        pix.save(out_path)
        rendered.append({
            "design_id": design_id,
            "pdf": os.path.basename(pdf_path),
            "page": pg["page"],
            "type": pg["type"],
            "image_path": out_path,
            "size": f"{pix.width}x{pix.height}",
        })

    doc.close()
    return rendered


# ════════════════════════════════════════════════════════════
# BATCH RUNNERS
# ════════════════════════════════════════════════════════════

def load_existing(output_dir: str, subdir: str, ext: str) -> set[str]:
    """Load set of already-processed Design IDs from output subdir."""
    existing = set()
    d = os.path.join(output_dir, subdir)
    if os.path.isdir(d):
        for root, dirs, files in os.walk(d):
            for f in files:
                if f.endswith(ext):
                    did = extract_design_id_from_path(f)
                    if did:
                        existing.add(did)
    return existing


def group_by_design(file_list: list[str], is_pdf: bool = False) -> tuple[dict, int]:
    """Group files by D-number. Returns (dict[did → [paths]], no_id_count)."""
    by_design: dict[str, list[str]] = {}
    no_id = 0
    for f in file_list:
        did = resolve_design_id(f, is_pdf=is_pdf)
        if did:
            by_design.setdefault(did, []).append(f)
        else:
            no_id += 1
    return by_design, no_id


def run_metadata_batch(scan_dir: str, output_dir: str, force: bool, dry_run: bool):
    """Path A: Extract metadata from all PDFs."""
    print(f"[METADATA] Scanning {scan_dir}...")
    all_files = scan_files(scan_dir, ".pdf")
    print(f"[METADATA] Found {len(all_files)} PDF files")

    by_design, no_id = group_by_design(all_files, is_pdf=True)
    print(f"[METADATA] {len(by_design)} unique designs ({no_id} without D-number)")

    if not force:
        # Check existing metadata
        meta_path = os.path.join(output_dir, "metadata", "designs.jsonl")
        existing = set()
        if os.path.exists(meta_path):
            with open(meta_path, encoding="utf-8") as f:
                for line in f:
                    if line.strip():
                        r = json.loads(line)
                        existing.add(r.get("design_id", ""))
        todo = {d: fs for d, fs in by_design.items() if d not in existing}
        print(f"[METADATA] Already extracted: {len(existing)}, Todo: {len(todo)}")
    else:
        todo = by_design

    if dry_run:
        print(f"[METADATA] Dry run — would extract {len(todo)} designs")
        return

    # Pick best PDF per design (largest = most complete techpack)
    meta_dir = os.path.join(output_dir, "metadata")
    os.makedirs(meta_dir, exist_ok=True)
    meta_path = os.path.join(meta_dir, "designs.jsonl")

    success = 0
    errors = 0

    with open(meta_path, "a", encoding="utf-8") as out_f:
        for did in sorted(todo):
            files = todo[did]
            # Pick largest PDF
            best = max(files, key=lambda p: os.path.getsize(p))
            try:
                meta = extract_pdf_metadata(best)
                if meta:
                    out_f.write(json.dumps(meta, ensure_ascii=False) + "\n")
                    success += 1
                else:
                    errors += 1
            except Exception as e:
                print(f"  ERROR {did}: {e}", file=sys.stderr)
                errors += 1

    print(f"[METADATA] Done: {success} extracted, {errors} errors → {meta_path}")


def run_pptx_batch(scan_dir: str, output_dir: str, force: bool, dry_run: bool):
    """Path B: Extract text from all PPTX files."""
    print(f"[PPTX] Scanning {scan_dir}...")
    all_files = scan_files(scan_dir, ".pptx")
    print(f"[PPTX] Found {len(all_files)} PPTX files")

    by_design, no_id = group_by_design(all_files)
    print(f"[PPTX] {len(by_design)} unique designs ({no_id} without D-number)")

    if not force:
        existing = load_existing(output_dir, "pptx", ".txt")
        todo = {d: fs for d, fs in by_design.items() if d not in existing}
        print(f"[PPTX] Already extracted: {len(existing)}, Todo: {len(todo)}")
    else:
        todo = by_design

    if dry_run:
        print(f"[PPTX] Dry run — would extract {len(todo)} designs")
        return

    pptx_out = os.path.join(output_dir, "pptx")
    os.makedirs(pptx_out, exist_ok=True)
    success = 0
    errors = 0

    for did in sorted(todo):
        for pptx_path in todo[did]:
            out_name = make_pptx_output_name(pptx_path, scan_dir)
            out_path = os.path.join(pptx_out, out_name)
            try:
                pptx_to_txt(pptx_path, out_path)
                success += 1
            except Exception as e:
                print(f"  ERROR {did} {os.path.basename(pptx_path)}: {e}", file=sys.stderr)
                errors += 1

    print(f"[PPTX] Done: {success} extracted, {errors} errors")


def run_pdf_batch(scan_dir: str, output_dir: str, force: bool, dry_run: bool):
    """Path C: Detect construction pages + render PNG from all PDFs."""
    print(f"[PDF→PNG] Scanning {scan_dir}...")
    all_files = scan_files(scan_dir, ".pdf")
    print(f"[PDF→PNG] Found {len(all_files)} PDF files")

    by_design, no_id = group_by_design(all_files, is_pdf=True)
    print(f"[PDF→PNG] {len(by_design)} unique designs ({no_id} without D-number)")

    if not force:
        existing = load_existing(output_dir, "pdf/callout_images", ".png")
        todo = {d: fs for d, fs in by_design.items() if d not in existing}
        print(f"[PDF→PNG] Already rendered: {len(existing)}, Todo: {len(todo)}")
    else:
        todo = by_design

    if dry_run:
        print(f"[PDF→PNG] Dry run — would process {len(todo)} designs")
        return

    img_dir = os.path.join(output_dir, "pdf", "callout_images")
    manifest_path = os.path.join(output_dir, "pdf", "callout_manifest.jsonl")
    os.makedirs(img_dir, exist_ok=True)

    total_designs = 0
    total_pages = 0
    no_callout = 0
    errors = 0

    with open(manifest_path, "a", encoding="utf-8") as manifest_f:
        for did in sorted(todo):
            best = max(todo[did], key=lambda p: os.path.getsize(p))
            try:
                pages = detect_construction_pages(best)
                if not pages:
                    no_callout += 1
                    continue
                rendered = render_pdf_pages(best, pages, img_dir, did)
                total_designs += 1
                total_pages += len(rendered)
                for r in rendered:
                    manifest_f.write(json.dumps(r, ensure_ascii=False) + "\n")
            except Exception as e:
                print(f"  ERROR {did}: {e}", file=sys.stderr)
                errors += 1

    print(f"[PDF→PNG] Done: {total_designs} designs / {total_pages} pages, "
          f"{no_callout} no callouts, {errors} errors")


# ════════════════════════════════════════════════════════════
# CLI
# ════════════════════════════════════════════════════════════

def main():
    p = argparse.ArgumentParser(
        description="Techpack 原始資料統一提取 — metadata + PPTX text + PDF construction PNG")
    p.add_argument("--scan-dir", required=True, help="掃描目錄")
    p.add_argument("--output-dir", default=None,
                   help="輸出目錄 (預設: star_schema/data/ingest/)")
    p.add_argument("--pptx-only", action="store_true", help="只跑 PPTX")
    p.add_argument("--pdf-only", action="store_true", help="只跑 PDF (metadata + PNG)")
    p.add_argument("--metadata-only", action="store_true", help="只跑 PDF metadata")
    p.add_argument("--force", action="store_true", help="忽略已處理，全部重跑")
    p.add_argument("--dry-run", action="store_true", help="列出待處理，不實際跑")
    args = p.parse_args()

    if args.output_dir is None:
        script_dir = Path(__file__).resolve().parent
        args.output_dir = str(script_dir.parent / "data" / "ingest")

    print(f"Output: {args.output_dir}")
    print(f"Scan:   {args.scan_dir}")
    print(f"Force: {args.force}, Dry run: {args.dry_run}")
    print()

    run_meta = not args.pptx_only
    run_pptx = not args.pdf_only and not args.metadata_only
    run_pdf = not args.pptx_only and not args.metadata_only

    if run_meta:
        run_metadata_batch(args.scan_dir, args.output_dir, args.force, args.dry_run)
        print()

    if run_pptx:
        run_pptx_batch(args.scan_dir, args.output_dir, args.force, args.dry_run)
        print()

    if run_pdf:
        run_pdf_batch(args.scan_dir, args.output_dir, args.force, args.dry_run)


if __name__ == "__main__":
    main()

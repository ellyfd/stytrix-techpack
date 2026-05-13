"""extract_pptx_all.py — Unified PPTX Techpack extractor (structured callout mode).

對 tp_samples_v2/<EIDH>/*.pptx 抽:
  - 整份 PPTX 文字 → pptx_text/<client>_<design>.txt (給 VLM / search 用)
  - 結構化 callout: 中文 zone + method + 推 ISO (主要產出)

2026-05-12 改版:
  ❌ 移除 PNG render (VLM pipeline 還沒接, 是 dead storage)
  ❌ 移除 per-slide schema (slides[]) - 改成 per-PPTX 一個 callout list
  ✅ 中文做工術語 → ISO inference (鎖鍊車→401, 平車→301, 拷克/地毯車→504, 三本車→406, 打結車→304, etc.)
  ✅ 中文 zone 偵測 (腰頭/口袋/領口/袖口/襠/側縫 等)
  ✅ Incremental --client X 模式 (保留其他 brand 既有資料)

用法:
  python scripts/extract_pptx_all.py [--limit N] [--workers N] [--client BRAND] [--reset]

需要:
  pip install python-pptx
"""
from __future__ import annotations
import argparse
import json
import re
import sys
import time
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path
from collections import Counter

SCRIPT_DIR = Path(__file__).resolve().parent
ROOT = SCRIPT_DIR.parent
TP_DIR = ROOT / "tp_samples_v2"
MANIFEST_PATH = ROOT / "m7_organized_v2" / "_fetch_manifest.csv"
OUT_DIR = ROOT / "outputs" / "extract"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_JSONL = OUT_DIR / "pptx_facets.jsonl"
PPTX_TEXT_DIR = OUT_DIR / "pptx_text"
PPTX_TEXT_DIR.mkdir(parents=True, exist_ok=True)
# PPTX_IMG_DIR removed 2026-05-12 — PNG render no longer done (VLM pipeline not wired)


# === Patterns for construction signal ===
ISO_RE = re.compile(r"\b(301|401|406|503|504|512|514|515|516|601|602|605|607)\b")
SEW_KW = [
    "COVERSTITCH", "OVERLOCK", "TOPSTITCH", "FLATLOCK", "FLATSEAM",
    "BARTACK", "BAR TACK", "BLINDHEM", "BLIND HEM", "EDGESTITCH",
    "BINDING", "SERGE", "CHAINSTITCH",
]


# ════════════════════════════════════════════════════════════
# 中文做工術語 → ISO inference (2026-05-12)
# 整理自 PPTX 翻譯文件實際用語 + IE 內部慣例
# ════════════════════════════════════════════════════════════

# 順序很重要: 長 keyword 先 (避免 "三本車" 被 "車" 搶走)
ZH_SEW_TO_ISO = [
    # === 雙針/三針 coverstitch family ===
    # 2026-05-12 補: 三本三針 / 三針三本 → 407 (Coverstitch 3-needle), 對齊 iso_dictionary 官方
    ("三本三針", "407"),     # 三本三針 = 3-needle coverstitch = 407
    ("三針三本", "407"),     # 變體
    ("三本針", "406"),
    ("三本車", "406"),
    ("三本", "406"),
    ("雙針包邊", "406"),
    ("覆蓋縫", "406"),

    # === 5-thread / 4-thread coverstitch ===
    ("五線", "516"),
    ("五條", "516"),

    # === Flatlock / Flatseam ===
    # 2026-05-12 補: 爬網 (605) / 併逢車 (607 變體寫法)
    # 注意衝突: iso_dictionary 把「併縫」標 602, IE memory 寫 607. 此處跟 IE 標準 (合 user 慣例).
    ("三針五線爬網", "605"),  # 605 官方 zh
    ("爬網", "605"),         # 605 alias
    ("併逢車", "607"),       # 607 官方 zh (注意是「逢」不是「縫」)
    ("併逢", "607"),
    ("平縫", "301"),
    ("FLATLOCK", "607"),
    ("FLAT SEAM", "602"),
    ("FLAT LOCK", "607"),

    # === Bartack / 套結 ===
    ("打結車", "304"),
    ("套結", "304"),
    ("補強車", "304"),
    ("BARTACK", "304"),
    ("BAR TACK", "304"),

    # === 鎖眼 (buttonhole, 304 zigzag) ===
    ("鎖眼", "304"),
    ("BUTTONHOLE", "304"),

    # === Overlock / Serger family ===
    ("拷克", "514"),
    ("拷邊", "514"),
    ("拷邊機", "514"),
    ("OVERLOCK", "504"),
    ("OVEREDGE", "504"),
    ("SERGE", "504"),
    ("SERGED", "504"),
    ("地毯車", "504"),

    # === 2-thread overedge (light) ===
    ("二條", "504"),
    ("二線", "504"),

    # === 4-thread overlock with safety stitch ===
    ("四條", "514"),
    ("四線安全", "514"),
    ("安全車", "504"),

    # === Chain stitch family ===
    # 2026-05-12 修: 鎖鏈/鎖鍊 alone → 401, 防 "壓單針鎖鏈" 被 "單針" 搶走判 301
    ("壓單針鎖鏈", "401"),    # 單針 chain stitch
    ("壓單針鎖鍊", "401"),
    ("單針鎖鏈", "401"),
    ("單針鎖鍊", "401"),
    ("鎖鍊車", "401"),
    ("鎖鏈車", "401"),
    ("鎖鏈", "401"),          # 鎖鏈 alone — chain stitch action (置於 單針 之前)
    ("鎖鍊", "401"),
    ("鏈式車", "401"),
    ("鏈式縫", "401"),
    ("鏈縫", "401"),
    ("鏈式縫紉", "401"),
    ("CHAINSTITCH", "401"),
    ("CHAIN STITCH", "401"),

    # === Double-needle ===
    ("雙針", "401"),

    # === Blind stitch ===
    ("暗縫", "301"),
    ("盲縫", "301"),
    ("BLIND STITCH", "301"),

    # === Lockstitch / single needle ===
    ("單針平車", "301"),
    ("單針", "301"),
    ("平車", "301"),
    ("壓單針", "301"),
    ("LOCKSTITCH", "301"),
    ("TOPSTITCH", "301"),

    # === 2026-05-12 補: 高頻 method keyword (audit 顯示 87% callout 無 ISO) ===
    ("壓三本", "406"),       # 壓三本車 = 三本 = 406
    ("併縫", "607"),         # 併縫 = Flatlock = 607 (IE 標準)
    ("反折壓線", "301"),     # 反折 + 壓線 = topstitch = 301
    ("反摺壓線", "301"),     # 同上 (繁體變體)
    ("反折兩次", "301"),     # 雙折 + 壓線
    ("反摺兩次", "301"),
    ("縫上口袋", "301"),     # 口袋縫合 default 301 (大宗)
    ("縫上", "301"),         # generic 縫合 (set in)
    ("剪接縫", "301"),       # 剪接 default 301 (panel seam, 大宗)
    ("剪接", "301"),         # generic 剪接
    ("壓線", "301"),         # generic topstitch (最高頻, 放最後 fallback)
    ("車死", "301"),         # 車死 = anchor stitch = 301 (lockstitch)
    ("埋邊", "514"),         # 埋邊 = overedge 變體
    ("貼邊壓線", "301"),     # 貼邊 + 壓線 = topstitch

    # === 滾邊 / 包邊 (binding, typically coverstitch 406) ===
    ("滾邊", "406"),
    ("包邊", "406"),
    ("BINDING", "406"),

    # === 直紋包邊 (Centric 8) ===
    ("直紋包邊", "406"),

    # === 簡體中文變體 (2026-05-12 加, 適用 TGT / Cat & Jack / 部分 ANF) ===
    ("锁链车", "401"), ("锁链", "401"),
    ("双针", "401"),
    ("压单针", "301"), ("单针平车", "301"), ("单针", "301"), ("平车", "301"),
    ("三本针", "406"), ("三本车", "406"),
    ("双针包边", "406"),
    ("滚边", "406"), ("包边", "406"),
    ("拷克", "514"), ("地毯车", "504"),
    ("打结车", "304"), ("套结", "304"), ("锁眼", "304"),
    ("链式车", "401"), ("链缝", "401"),
    ("平缝", "301"),
    ("暗缝", "301"), ("盲缝", "301"),
]

ZH_SEW_TO_ISO_COMPILED = [(re.compile(re.escape(kw), re.IGNORECASE), iso) for kw, iso in ZH_SEW_TO_ISO]


# === 2026-05-12 加 / 修: Negative filter — 過濾「不是做工方法」的內容 ===
# 這些 text 即使含 sewing-like keyword 也不該推 ISO (避免 false-positive)
# audit 樣本顯示: "車線:TEX 27" / "洗標#LBLST1033" / "圖#3" / "請改善" 等
#
# ❗ 拿掉了「距 XX 縫 3/8"」「完成下擺上」兩條 pattern (2026-05-12 第二輪修):
#    「距袋口邊3/8"壓單針平車(301)」是合法 sewing instruction 含位置描述,
#    被原 pattern 誤殺,讓 ISO 推導率漲幅從預期 +17pp 縮到 +5pp.
#    現在只留下「整段就是 non-method」的高精度 pattern.
NON_METHOD_PATTERNS = [
    re.compile(r"^車線\s*[:：]"),        # 車線: TEX 27 → 規格不是方法
    re.compile(r"^針密\s*[:：]"),        # 針密: 11 SPI
    re.compile(r"^線材\s*[:：]"),        # 線材規格
    re.compile(r"洗標\s*[#＃\-﹣]"),     # 洗標#LBLST1033 / 洗標-XXX → 標籤名
    re.compile(r"^洗標"),                # 洗標 在開頭 (位置標註)
    re.compile(r"^圖\s*[#＃]"),          # 圖#3 → 圖號參照
    re.compile(r"^圖[號片]"),            # 圖號 / 圖片
    re.compile(r"請(改善|修正|確認|檢查|注意)"),  # QC 評語
    re.compile(r"反黃尺寸"),             # 反黃尺寸 → 尺寸標註
    re.compile(r"^尺寸表"),              # 尺寸表
]


def _is_non_method(text: str) -> bool:
    """判斷一段 text 是不是 sewing method (return True 表示『不是做工方法』).

    重要設計: 只有當整段 text 就是 non-method 時才回 True.
    含「距 XX 邊 3/8"」這種 dimension 描述但同時有 sewing keyword 的合法 instruction 不算 non-method.
    """
    if not text:
        return True
    for pat in NON_METHOD_PATTERNS:
        if pat.search(text):
            return True
    return False


# 中文 zone (部位) 偵測
# 順序: 長 keyword 優先 (避免 "袖" 搶走 "袖口")
ZH_ZONES = [
    # Tops
    ("領口", "領口"), ("領圍", "領圍"), ("領片", "領片"), ("領座", "領座"),
    ("肩縫", "肩縫"), ("肩線", "肩線"), ("肩帶", "肩帶"),
    ("袖口", "袖口"), ("袖籠", "袖籠"), ("袖窿", "袖籠"),
    ("袖克夫", "袖克夫"), ("袖標", "袖標"),
    ("胸口", "胸"), ("前胸", "前胸"),
    ("下擺", "下擺"), ("底擺", "下擺"),
    ("肩", "肩"),  # last fallback
    ("袖", "袖"),  # last fallback for sleeve

    # Front / back / center
    ("前中", "前中"), ("後中", "後中"),
    ("前片", "前片"), ("後片", "後片"),
    ("前領", "前領"), ("後領", "後領"),
    ("領", "領"),  # generic neck fallback

    # 簡體中文 zone variants
    ("腰头", "腰頭"), ("腰绳", "腰繩"), ("腰带", "腰帶"),
    ("袖笼", "袖籠"), ("袖窿", "袖籠"),
    ("领口", "領口"), ("领围", "領圍"), ("领片", "領片"),
    ("肩缝", "肩縫"),
    ("下摆", "下擺"), ("底摆", "下擺"),
    ("前裆", "前襠"), ("后裆", "後襠"), ("内裆", "內襠"), ("裆", "襠"),
    ("裤口", "褲口"),
    ("脅边", "側縫"), ("侧缝", "側縫"), ("侧边", "側縫"),
    ("剪接线", "剪接"),
    ("门襟", "門襟"),
    ("开岔", "開岔"), ("开叉", "開岔"),
    ("反折", "反摺"),

    # Bottoms (繁體)
    ("腰頭", "腰頭"), ("腰繩", "腰繩"), ("腰帶", "腰帶"),
    ("腰圍", "腰圍"), ("腰", "腰"),
    ("臀圍", "臀圍"), ("臀", "臀"),
    ("前襠", "前襠"), ("後襠", "後襠"), ("內襠", "內襠"),
    ("襠底", "襠底"), ("襠", "襠"),
    ("褲口", "褲口"), ("腳口", "褲口"),
    ("褲檔", "褲檔"),
    ("膝", "膝"),

    # Seams
    ("側縫", "側縫"), ("脅邊", "側縫"), ("側邊", "側縫"),
    ("剪接線", "剪接"), ("剪接", "剪接"),

    # Pockets / closures
    ("口袋", "口袋"),
    ("拉鍊", "拉鍊"), ("拉鏈", "拉鍊"),
    ("門襟", "門襟"),
    ("開岔", "開岔"), ("開叉", "開岔"),
    ("鎖眼", "鎖眼"),  # also action keyword but also a location
    ("反摺", "反摺"), ("反折", "反摺"),

    # Other / accessories
    ("罩杯", "罩杯"),
    ("LOGO", "LOGO"),
]

ZH_ZONES_COMPILED = [(re.compile(re.escape(kw)), name) for kw, name in ZH_ZONES]


# === 中文 zone → 五階 L1 code mapping (38 official L1) ===
# 對應 stytrix-techpack/l2_l3_ie/<L1>.json 38 個官方部位
# 中文 zone keyword (來自 ZH_ZONES) → L1 code (大寫 2 字母)
ZH_ZONE_TO_L1 = {
    # Tops
    "領": "NK", "領口": "NK", "領圍": "NK", "領片": "NK", "領座": "NK",
    "前領": "NK", "後領": "NK",
    "領襟": "NP",
    "領貼條": "NT",
    "肩": "SH", "肩縫": "SH", "肩線": "SH",
    "肩帶": "ST",
    "袖孔": "AE",  # 袖孔 = 無袖款式 armhole (official)
    "袖": "AH",  # 2026-05-12 補: 袖 (generic) → AH 袖圍 (audit 顯示 57,412 件無 L1)
    "袖籠": "AH", "袖窿": "AH", "袖圍": "AH",  # armhole edge
    "袖口": "SL", "袖克夫": "SL",
    "袖叉": "SP",
    "袖標": "LB",

    # Body
    # 胸/前胸 不 mapping (location 不是 part)
    "下擺": "BM", "底擺": "BM",
    # 前中/後中/前片/後片 不直接 mapping (是 location 不是 part, 由 method 內容決定 L1)
    "剪接": "SA",  # 上身剪接 → SA (剪接線_上身類)

    # Bottoms
    "腰頭": "WB", "腰繩": "DC", "腰帶": "WB", "腰圍": "WB", "腰": "WB",
    "臀": "PS", "臀圍": "PS",
    "前襠": "RS", "後襠": "RS", "內襠": "RS", "襠底": "RS", "襠": "RS", "褲檔": "RS",
    "褲口": "LO", "腳口": "LO",
    "膝": "PS",
    "褲合身": "PS",
    "裙合身": "SR",

    # Seams
    "側縫": "SS", "脅邊": "SS",

    # Pockets / closures
    "口袋": "PK",
    "袋蓋": "FP",
    "拉鍊": "ZP",
    "門襟": "PL",
    "前立": "FY",
    "開岔": "BP",  # 襬叉 (bottom vent/slit)

    # Accessories / decoration
    "罩杯": "OT",
    "LOGO": "LB",
    "商標": "LB",
    "貼合": "BN",
    "釦鎖": "BS",
    "繩": "DC", "繩類": "DC",
    "裝飾片": "DP",
    "帽子": "HD",
    "釦環": "HL",
    "鎖眼": "HL",  # 2026-05-12 補: 鎖眼 (buttonhole) → HL 釦環 (NOT KH, KH 是極少見的鑰匙孔狀領口)
    "裡布": "LI",
    "帶絆": "LP",
    "褶": "PD",
    "行縫": "QT",
    "拇指洞": "TH",
    # 反摺 是動作不是部位, 不 mapping (context-dependent: 腰頭反折/褲口反折/袖口反折)
    # 前中/後中/前胸/前片/後片/胸 — 官方 38 L1 無對應碼, 是 piece boundary 非 detail zone
}

# === Load official L1 + ISO canonical names from stytrix-techpack repo ===
# Source of truth: stytrix-techpack/data/runtime/l1_standard_38.json (38 L1) +
#                  stytrix-techpack/data/runtime/iso_dictionary.json (ISO 中英文)

def _load_l1_dict() -> dict[str, dict]:
    """Load l1_standard_38.json {code → {zh, en}}. Fallback hardcoded if not found."""
    candidates = [
        Path(r"C:\temp\stytrix-techpack\data\runtime\l1_standard_38.json"),
        Path(__file__).resolve().parent.parent.parent.parent / "stytrix-techpack" / "data" / "runtime" / "l1_standard_38.json",
        SCRIPT_DIR.parent.parent / "stytrix-techpack" / "data" / "runtime" / "l1_standard_38.json",
    ]
    for c in candidates:
        if c.exists():
            try:
                with open(c, encoding="utf-8") as f:
                    data = json.load(f)
                return data.get("codes", {})
            except Exception:
                continue
    # Fallback hardcoded (2026-05-12 snapshot of official)
    return {
        "AE": {"zh": "袖孔", "en": "Armhole"},
        "AH": {"zh": "袖圍", "en": "Armhole"},
        "BM": {"zh": "下襬", "en": "Bottom"},
        "BN": {"zh": "貼合", "en": "Bonding"},
        "BP": {"zh": "襬叉", "en": "Bottom Placket"},
        "BS": {"zh": "釦鎖", "en": "Botton/Snap"},
        "DC": {"zh": "繩類", "en": "Drawcord"},
        "DP": {"zh": "裝飾片", "en": "Decoration Pieces"},
        "FP": {"zh": "袋蓋", "en": "Flap"},
        "FY": {"zh": "前立", "en": "Front Fly"},
        "HD": {"zh": "帽子", "en": "Hoodie"},
        "HL": {"zh": "釦環", "en": "Hook and Loop"},
        "KH": {"zh": "Keyhole", "en": "Keyhole"},
        "LB": {"zh": "商標", "en": "Lable"},
        "LI": {"zh": "裡布", "en": "Lining"},
        "LO": {"zh": "褲口", "en": "Leg Opening"},
        "LP": {"zh": "帶絆", "en": "Loop"},
        "NK": {"zh": "領", "en": "Neck"},
        "NP": {"zh": "領襟", "en": "Neck Placket"},
        "NT": {"zh": "領貼條", "en": "Neck Tape"},
        "OT": {"zh": "其它", "en": "Other"},
        "PD": {"zh": "褶", "en": "Pleats/Darts"},
        "PK": {"zh": "口袋", "en": "Pocket"},
        "PL": {"zh": "門襟", "en": "Placket"},
        "PS": {"zh": "褲合身", "en": "Pants Side Seam"},
        "QT": {"zh": "行縫(固定棉)", "en": "Quilted"},
        "RS": {"zh": "褲襠", "en": "Rise"},
        "SA": {"zh": "剪接線_上身類", "en": "Seam/Trim/Piecing/Panel Top"},
        "SB": {"zh": "剪接線_下身類", "en": "Seam/Trim/Piecing/Panel Pants"},
        "SH": {"zh": "肩", "en": "Shoulder"},
        "SL": {"zh": "袖口", "en": "Sleeve"},
        "SP": {"zh": "袖叉", "en": "Sleeve Placket"},
        "SR": {"zh": "裙合身", "en": "Skirt Side Seam"},
        "SS": {"zh": "脅邊", "en": "Side Seam"},
        "ST": {"zh": "肩帶", "en": "Strap"},
        "TH": {"zh": "拇指洞", "en": "Thumb Hole"},
        "WB": {"zh": "腰頭", "en": "Waistband"},
        "ZP": {"zh": "拉鍊", "en": "Zipper"},
    }


L1_DICT = _load_l1_dict()
# Convenience accessor
L1_CODE_TO_ZH_NAME = {code: info["zh"] for code, info in L1_DICT.items()}
L1_CODE_TO_EN_NAME = {code: info["en"] for code, info in L1_DICT.items()}


def _zone_to_l1(zone: str) -> str | None:
    """中文 zone → L1 2-letter code (None if no mapping)."""
    return ZH_ZONE_TO_L1.get(zone)


# 2026-05-12 加: 顯式括號 ISO 抽取 — 「壓單針平車(301)」這類直接含 ISO code 的 callout
# 最高 confidence: 翻譯文件 + IE 寫法都會在 method 後加 (NNN) 顯式標註
# 對齊 stytrix-techpack/data/runtime/iso_dictionary.json 14 codes (含複合碼 514+401 / 514+605)
ISO_BRACKET_RE = re.compile(
    r"[(（]\s*("
    r"514\+401|514\+605|"  # 複合碼優先 (避免被 514 / 401 / 605 拆成兩個)
    r"301|304|401|406|407|503|504|512|514|516|602|605|607"
    r")\s*[)）]"
)


def _infer_iso_from_zh(text):
    """從中文做工描述 infer ISO code list.

    Priority 順序:
      1. 顯式括號 ISO: "壓單針平車(301)" → 取 301 (最高 confidence)
      2. Keyword inference: "壓單針平車" → 301 (heuristic)

    2026-05-12 加 negative filter — 含「車線:」「洗標#」「圖#」「請改善」等非做工內容 → 不推 ISO,
    避免「車線: TEX 27」被誤判成「車」→ 301
    """
    if _is_non_method(text):
        return []
    isos = []
    seen = set()
    # Priority 1: explicit (NNN) brackets — highest confidence
    for m in ISO_BRACKET_RE.finditer(text):
        iso = m.group(1)
        if iso not in seen:
            isos.append(iso)
            seen.add(iso)
    # Priority 2: keyword-based heuristic inference
    for pat, iso in ZH_SEW_TO_ISO_COMPILED:
        if pat.search(text) and iso not in seen:
            isos.append(iso)
            seen.add(iso)
    return isos


def _extract_zones(text):
    """從 text scan zone keyword (順序敏感, 長詞優先)."""
    zones = []
    seen = set()
    # consume text greedy: longer keywords first
    for pat, name in ZH_ZONES_COMPILED:
        if pat.search(text) and name not in seen:
            zones.append(name)
            seen.add(name)
    return zones


def _parse_slide_constructions(slide_text, slide_num):
    """從一張 slide 的 text 拆 callout 結構.

    Strategy:
      Pattern 1 (best): "<zone>: <description>" — 整段 zone + descriptions block
      Pattern 2 (fallback): line-by-line — 每行 scan zone + iso, 有任一 match 就當 callout
    """
    constructions = []
    lines = [l.strip() for l in slide_text.split("\n")]

    # === Pattern 1: 找 "<zone>:" 起頭, 抓接下來幾行直到下一個 zone 開頭 ===
    blocks = []  # [(zone_label, [description_lines])]
    current_zone = None
    current_block = []
    zone_colon_re = re.compile(r"^(.{1,15}?):\s*$|^(.{1,15}?):\s+(.+)")

    for line in lines:
        if not line:
            if current_zone and current_block:
                blocks.append((current_zone, current_block))
                current_zone, current_block = None, []
            continue
        # Detect zone-colon line
        m = zone_colon_re.match(line)
        if m:
            zone_text = m.group(1) or m.group(2) or ""
            inline_rest = m.group(3) if m.lastindex == 3 else ""
            # Check if zone_text matches any known zone
            zones_found = _extract_zones(zone_text)
            if zones_found:
                # Save previous block
                if current_zone and current_block:
                    blocks.append((current_zone, current_block))
                current_zone = zones_found[0]
                current_block = []
                if inline_rest:
                    current_block.append(inline_rest)
                continue
        # Regular description line
        if current_zone:
            # Skip noise (boilerplate)
            if "CONFIDENTIAL" in line.upper() or "ALL RIGHTS RESERVED" in line.upper():
                continue
            current_block.append(line)

    if current_zone and current_block:
        blocks.append((current_zone, current_block))

    # === Convert blocks into callouts ===
    for zone, descs in blocks:
        for desc in descs:
            # Strip leading "-" / "•" / "*"
            d = desc.lstrip("-•*").strip()
            if not d or len(d) < 2:
                continue
            # 排掉純數字/英文 noise (< 3 chars 但含中文 OK)
            if len(d) < 3 and not any("\u4e00" <= ch <= "\u9fff" for ch in d):
                continue
            isos = _infer_iso_from_zh(d)
            callout = {
                "zone": zone,
                "method": d[:300],
                "_source_slide": slide_num,
            }
            # Map zone to L1 (五階官方代號)
            l1 = _zone_to_l1(zone)
            # 2026-05-12 加: zone 沒 L1 mapping 時 (反摺/前中/後中/前胸/前片/後片/胸),
            # 從 method 內容反向 scan body-part keyword 推 L1 — 救回 ~30K callouts
            if not l1:
                method_zones = _extract_zones(d)
                for mz in method_zones:
                    candidate = _zone_to_l1(mz)
                    if candidate:
                        l1 = candidate
                        callout["L1_inferred_from"] = mz  # audit trail
                        break
            if l1:
                callout["L1"] = l1
                callout["L1_name"] = L1_CODE_TO_ZH_NAME.get(l1, "")
            if isos:
                callout["iso"] = isos[0]
                if len(isos) > 1:
                    callout["iso_alt"] = isos[1:]
            constructions.append(callout)

    # === Pattern 2 fallback: if no blocks found, scan each line independently ===
    if not callouts:
        for line in lines:
            # Allow short Chinese lines (e.g. "暗縫")
            if not line or (len(line) < 5 and not any("\u4e00" <= ch <= "\u9fff" for ch in line)):
                continue
            if "CONFIDENTIAL" in line.upper() or "ALL RIGHTS RESERVED" in line.upper():
                continue
            zones = _extract_zones(line)
            isos = _infer_iso_from_zh(line)
            if zones or isos:
                callout = {
                    "method": line[:300],
                    "_source_slide": slide_num,
                }
                if zones:
                    callout["zone"] = zones[0]
                    if len(zones) > 1:
                        callout["zone_alt"] = zones[1:]
                    # Map first zone to L1
                    l1 = _zone_to_l1(zones[0])
                    if l1:
                        callout["L1"] = l1
                        callout["L1_name"] = L1_CODE_TO_ZH_NAME.get(l1, "")
                if isos:
                    callout["iso"] = isos[0]
                    if len(isos) > 1:
                        callout["iso_alt"] = isos[1:]
                constructions.append(callout)

    return constructions



# Folder name parsing (same as extract_pdf_all)
KNOWN_CLIENT_TOKENS = [
    "DICKS_SPORTING_GOODS", "ABERCROMBIE_AND_FITCH",
    "OLD_NAVY", "GAP_OUTLET", "BANANA_REPUBLIC", "WAL-MART-CA",
    "GAP", "DICKS", "ATHLETA", "UNDER_ARMOUR", "KOHLS", "A_AND_F", "GU", "BEYOND_YOGA",
    "HIGH_LIFE_LLC", "WAL-MART", "QUINCE", "HALARA", "NET",
    "JOE_FRESH", "BRFS", "SANMAR", "DISTANCE", "ZARA",
    "ASICS-EU", "TARGET", "LEVIS", "CATO", "SMART_CLOTHING",
]

CLIENT_RAW_TO_CODE = {
    "OLD NAVY": "ONY", "TARGET": "TGT", "GAP": "GAP", "GAP OUTLET": "GAP",
    "DICKS SPORTING GOODS": "DKS", "DICKS": "DKS", "ATHLETA": "ATH",
    "UNDER ARMOUR": "UA", "KOHLS": "KOH", "A AND F": "ANF", "A & F": "ANF",
    "ABERCROMBIE & FITCH": "ANF", "ABERCROMBIE AND FITCH": "ANF",
    "GU": "GU", "BEYOND YOGA": "BY", "HIGH LIFE LLC": "HLF", "WAL-MART": "WMT",
    "WAL-MART-CA": "WMT", "QUINCE": "QCE", "HALARA": "HLA", "NET": "NET",
    "JOE FRESH": "JF", "BANANA REPUBLIC": "BR", "BRFS": "BR", "SANMAR": "SAN",
    "DISTANCE": "DST", "ZARA": "ZAR", "ASICS-EU": "ASICS", "LEVIS": "LEV",
    "CATO": "CATO", "SMART CLOTHING": "SMC",
}


def _load_manifest_lookup() -> dict:
    import csv
    lookup = {}
    if not MANIFEST_PATH.exists():
        return lookup
    with open(MANIFEST_PATH, "r", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            eidh = (row.get("Eidh") or "").strip()
            if eidh:
                lookup[eidh] = {
                    "客戶": (row.get("客戶") or "").strip(),
                    "報價款號": (row.get("報價款號") or "").strip(),
                    "Item": (row.get("Item") or "").strip(),
                    "HEADER_SN": (row.get("HEADER_SN") or "").strip(),
                }
    return lookup


_MANIFEST_LOOKUP = {}


def _init_worker(manifest_lookup):
    global _MANIFEST_LOOKUP
    _MANIFEST_LOOKUP = manifest_lookup


def _parse_folder_name(folder_name: str, manifest_lookup: dict = None) -> dict:
    lookup = manifest_lookup if manifest_lookup is not None else _MANIFEST_LOOKUP
    parts = folder_name.split("_", 1)
    eidh = parts[0] if parts else None
    design_suffix = parts[1] if len(parts) > 1 else ""
    info = lookup.get(eidh, {})
    client_raw = info.get("客戶", "")
    design_id = info.get("報價款號") or design_suffix
    client_code = CLIENT_RAW_TO_CODE.get(client_raw.upper().strip(), client_raw[:6].upper() if client_raw else "UNKNOWN")
    return {
        "eidh": eidh,
        "hsn": info.get("HEADER_SN", ""),
        "client_raw": client_raw,
        "client_code": client_code,
        "design_id": design_id,
        "item": info.get("Item", ""),
    }


def _slide_score(text: str, n_images: int) -> int:
    """score per slide for construction signal (含中文 keyword 2026-05-12 加)."""
    score = 0
    upper = text.upper()
    if ISO_RE.search(text): score += 3
    if any(k in upper for k in SEW_KW): score += 2
    if n_images >= 1: score += 1
    if "CALLOUT" in upper or "CONSTRUCTION" in upper: score += 2
    # 中文 sewing keyword (聚陽自家 PPT 95% 中文)
    ZH_SEW_TOKENS = ["鎖鍊", "锁链", "平車", "平车", "三本", "拷克",
                     "打結", "打结", "鎖眼", "锁眼", "雙針", "双针",
                     "暗縫", "暗缝", "覆蓋", "包邊", "包边", "滾邊", "滚边",
                     "壓線", "压线", "壓1/", "壓3/", "做工", "車縫"]
    if any(t in text for t in ZH_SEW_TOKENS): score += 2
    # 中文 zone keyword (suggesting construction annotations)
    ZH_ZONE_TOKENS = ["腰頭", "腰头", "領口", "领口", "袖口", "肩縫", "肩缝",
                      "口袋", "下擺", "下摆", "門襟", "门襟", "前襠", "前裆"]
    if any(t in text for t in ZH_ZONE_TOKENS): score += 1
    return score


def _worker_extract(args_tuple) -> dict:
    folder_path_str, _unused = args_tuple  # do_render flag retained for backward compat but no longer used
    folder = Path(folder_path_str)
    meta = _parse_folder_name(folder.name)
    eidh = meta["eidh"]
    client_code = meta["client_code"]
    design_id = meta["design_id"]

    pptxs = sorted(folder.glob("*.pptx"))
    if not pptxs:
        return {"eidh": eidh, "client_code": client_code, "design_id": design_id, "_status": "no_pptx"}

    facets = {
        "eidh": eidh,
        "client_code": client_code,
        "client_raw": meta["client_raw"],
        "design_id": design_id,
        "pptx_files": [],
        "n_slides_total": 0,
        "n_construction_slides": 0,
        "constructions": [],           # 2026-05-12: 結構化 callout list (per-PPTX flat)
        "n_constructions": 0,
        "raw_text_file": None,
        "_status": "ok",
    }

    try:
        from pptx import Presentation
    except ImportError:
        return {**facets, "_status": "no_python_pptx"}

    all_text_parts = []

    for pptx in pptxs:
        facets["pptx_files"].append(pptx.name)
        try:
            prs = Presentation(str(pptx))
        except Exception:
            continue

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            n_images = 0
            for shape in slide.shapes:
                if hasattr(shape, "image") or shape.shape_type == 13:  # picture
                    n_images += 1
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)
            slide_text = "\n".join(slide_text_parts).strip()
            score = _slide_score(slide_text, n_images)

            facets["n_slides_total"] += 1
            if score >= 3:
                facets["n_construction_slides"] += 1

            # 結構化 callout 抽取 (中文 zone + method + ISO)
            slide_constructions = _parse_slide_constructions(slide_text, slide_idx)
            for c in slide_callouts:
                c["_source_pptx"] = pptx.name
            facets["constructions"].extend(slide_callouts)

            # Raw text dump for VLM / search (留著, 不太大)
            all_text_parts.append(f"\n========== {pptx.name} slide {slide_idx} (score={score}) ==========\n")
            all_text_parts.append(slide_text)

    facets["n_constructions"] = len(facets["constructions"])

    # write combined raw text
    if all_text_parts:
        txt_path = PPTX_TEXT_DIR / f"{client_code}_{design_id}.txt"
        try:
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write("\n".join(all_text_parts))
            facets["raw_text_file"] = txt_path.relative_to(ROOT).as_posix()
        except Exception:
            pass

    return facets


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--limit", type=int, default=0)
    p.add_argument("--workers", type=int, default=4)
    p.add_argument("--client", help="只跑特定 brand code")
    p.add_argument("--no-render", action="store_true", help="(deprecated, PNG render 已永久移除)")
    p.add_argument("--reset", action="store_true",
                   help="--client 模式時清掉整個 jsonl 重抽 (預設保留其他 brand 的舊資料)")
    args = p.parse_args()

    manifest_lookup = _load_manifest_lookup()
    print(f"[manifest] loaded {len(manifest_lookup):,} EIDH lookup entries")
    global _MANIFEST_LOOKUP
    _MANIFEST_LOOKUP = manifest_lookup

    folders = sorted(d for d in TP_DIR.iterdir() if d.is_dir())
    if args.client:
        folders = [f for f in folders if _parse_folder_name(f.name, manifest_lookup)["client_code"] == args.client]
    if args.limit:
        folders = folders[:args.limit]
    print(f"[scan] {TP_DIR}: {len(folders)} EIDH folders (client={args.client or 'all'})")

    # === Incremental mode: --client X 時保留其他 brand 的舊 entries ===
    preserved_lines = []
    if args.client and not args.reset and OUT_JSONL.exists():
        n_dropped = 0
        n_kept = 0
        try:
            with open(OUT_JSONL, "r", encoding="utf-8") as fin:
                for line in fin:
                    try:
                        d = json.loads(line)
                    except json.JSONDecodeError:
                        continue
                    if d.get("client_code") == args.client:
                        n_dropped += 1
                    else:
                        preserved_lines.append(line.rstrip("\n"))
                        n_kept += 1
            print(f"[incremental] keeping {n_kept:,} existing rows (other brands), "
                  f"dropping {n_dropped:,} old rows for client={args.client}")
        except Exception as e:
            print(f"  [!] incremental read failed: {e}, falling back to fresh extract", file=sys.stderr)
            preserved_lines = []

    t0 = time.time()
    stats = Counter()
    by_client = Counter()
    by_client_with_construction = Counter()
    total_slides = 0
    total_construction_slides = 0
    total_callouts = 0

    with open(OUT_JSONL, "w", encoding="utf-8") as fout, \
         ProcessPoolExecutor(max_workers=args.workers,
                             initializer=_init_worker,
                             initargs=(manifest_lookup,)) as ex:
        # 先寫回保留的其他 brand entries
        for line in preserved_lines:
            fout.write(line + "\n")
        do_render = False  # 永遠不 render PNG (2026-05-12)
        futures = {ex.submit(_worker_extract, (str(d), do_render)): d.name for d in folders}
        for i, fut in enumerate(as_completed(futures)):
            try:
                r = fut.result()
            except Exception as e:
                print(f"  [!] {futures[fut]}: {e}", file=sys.stderr)
                stats["worker_err"] += 1
                continue
            status = r.get("_status", "?")
            stats[status] += 1
            total_callouts += r.get("n_constructions", 0)
            cl = r.get("client_code", "UNKNOWN")
            by_client[cl] += 1
            if r.get("n_construction_slides", 0) > 0:
                by_client_with_construction[cl] += 1
            total_slides += r.get("n_slides_total", 0)
            total_construction_slides += r.get("n_construction_slides", 0)
            fout.write(json.dumps(r, ensure_ascii=False) + "\n")
            if (i + 1) % 50 == 0:
                rate = (i + 1) / max(time.time() - t0, 0.1)
                eta = (len(folders) - i - 1) / rate / 60
                print(f"  [{i+1}/{len(folders)}] rate={rate:.1f}/s ETA={eta:.0f}min", flush=True)

    elapsed_min = (time.time() - t0) / 60
    print(f"\n[done] {sum(stats.values())} folders in {elapsed_min:.1f} min")
    print(f"  total slides: {total_slides:,}, construction slides (score>=3): {total_construction_slides:,}")
    print(f"\nstatus:")
    for s, n in stats.most_common():
        print(f"  {s:<15} {n:>6}")
    print(f"\nby client (total / w/construction):")
    for cl in sorted(by_client.keys(), key=lambda x: -by_client[x]):
        print(f"  {cl:<8} {by_client[cl]:>5} {by_client_with_construction[cl]:>5}")
    print(f"\noutput: {OUT_JSONL}")
    print(f"raw text: {PPTX_TEXT_DIR}")


if __name__ == "__main__":
    main()
